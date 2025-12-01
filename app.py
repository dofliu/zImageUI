import torch
import os
from diffusers import ZImagePipeline
from flask import Flask, render_template, request, jsonify, send_from_directory, send_file
from datetime import datetime
import base64
from io import BytesIO
import random
import config  # 導入配置檔案
import json
import zipfile
import tempfile
from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

app = Flask(__name__)

# 從配置檔案讀取路徑設定
cache_path = config.CACHE_PATH
output_path = config.OUTPUT_PATH
history_file = os.path.join(output_path, "history.json")
templates_file = os.path.join(os.path.dirname(__file__), "templates.json")
os.makedirs(cache_path, exist_ok=True)
os.makedirs(output_path, exist_ok=True)

# 全域變數存放模型
pipe = None

# 歷史記錄管理
def load_history():
    """載入歷史記錄"""
    if os.path.exists(history_file):
        try:
            with open(history_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return []
    return []

def save_history(history):
    """儲存歷史記錄"""
    try:
        with open(history_file, 'w', encoding='utf-8') as f:
            json.dump(history, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"儲存歷史記錄失敗: {e}")

def add_to_history(prompt, filename, tags=None):
    """新增歷史記錄"""
    history = load_history()
    history_item = {
        'id': f"{int(datetime.now().timestamp() * 1000)}_{random.randint(1000, 9999)}",
        'prompt': prompt,
        'filename': filename,
        'timestamp': datetime.now().isoformat(),
        'image_url': f'/images/{filename}',
        'tags': tags if tags else []
    }
    history.insert(0, history_item)  # 最新的在前面
    # 限制歷史記錄數量 (最多50筆)
    if len(history) > 50:
        history = history[:50]
    save_history(history)
    return history_item

def initialize_model():
    """初始化模型 (只執行一次)"""
    global pipe
    if pipe is None:
        # 檢查本地快取是否存在
        model_cache_exists = os.path.exists(os.path.join(cache_path, "models--Tongyi-MAI--Z-Image-Turbo"))
        if model_cache_exists:
            print("✓ 發現本地快取,從硬碟載入模型...")
        else:
            print("✗ 未發現本地快取,將從 Hugging Face 下載模型 (這需要較長時間)...")

        import time
        start_time = time.time()

        pipe = ZImagePipeline.from_pretrained(
            "Tongyi-MAI/Z-Image-Turbo",
            torch_dtype=torch.bfloat16,
            low_cpu_mem_usage=True,
            cache_dir=cache_path,
            use_safetensors=True,
            local_files_only=True,  # 允許從網路下載,但會優先使用本地快取
        )

        # 針對 12GB VRAM 的優化設定
        print("⟳ 啟用 Sequential CPU Offload (更激進的顯存優化)...")
        pipe.enable_sequential_cpu_offload()
        # 啟用額外的 VRAM 優化
        optimizations = []

        if config.ENABLE_ATTENTION_SLICING:
            if hasattr(pipe, 'enable_attention_slicing'):
                try:
                    pipe.enable_attention_slicing("auto")
                    optimizations.append("Attention Slicing")
                except Exception as e:
                    print(f"! Attention Slicing 啟用失敗: {e}")

        if config.ENABLE_VAE_SLICING:
            if hasattr(pipe, 'enable_vae_slicing'):
                try:
                    pipe.enable_vae_slicing()
                    optimizations.append("VAE Slicing")
                except Exception as e:
                    print(f"! VAE Slicing 啟用失敗: {e}")

        if config.ENABLE_XFORMERS:
            if hasattr(pipe, 'enable_xformers_memory_efficient_attention'):
                try:
                    pipe.enable_xformers_memory_efficient_attention()
                    optimizations.append("xFormers Attention")
                except Exception as e:
                    print(f"! xFormers 啟用失敗: {e}")

        if optimizations:
            print(f"✓ 已啟用額外優化: {', '.join(optimizations)}")

        # 嘗試啟用 VAE Tiling (如果 pipeline 支援)
        try:
            if hasattr(pipe, 'enable_vae_tiling'):
                pipe.enable_vae_tiling()
                print("✓ 已啟用 VAE Tiling (優化高解析度生成)")
            else:
                print("! ZImagePipeline 不支援 VAE Tiling，跳過此優化")
        except Exception as e:
            print(f"! VAE Tiling 啟用失敗: {e}")

        # 嘗試啟用 Flash Attention 加速 (如果環境支援)
        try:
            if hasattr(pipe, 'transformer') and hasattr(pipe.transformer, 'set_attention_backend'):
                pipe.transformer.set_attention_backend("flash")
                print("✓ 已啟用 Flash Attention 加速")
            else:
                print("! Flash Attention 不可用，使用預設 Attention")
        except Exception as e:
            print(f"! Flash Attention 啟用失敗: {e}")

        # 注意: 已使用 enable_model_cpu_offload()，不再需要 pipe.to("cuda")

        elapsed_time = time.time() - start_time
        print(f"✓ 模型載入完成! (耗時 {elapsed_time:.1f} 秒)")
    else:
        print("✓ 模型已在記憶體中,跳過載入")

@app.route('/')
def index():
    """首頁"""
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_image():
    """生成圖片 API"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', '')
        style_keywords = data.get('style_keywords', '')  # 風格關鍵字
        custom_width = data.get('width')  # 自定義寬度
        custom_height = data.get('height')  # 自定義高度

        if not prompt:
            return jsonify({'error': '請輸入提示詞'}), 400

        # 組合風格關鍵字到提示詞
        if style_keywords:
            full_prompt = f"{prompt}, {style_keywords}"
        else:
            full_prompt = prompt

        # 確定使用的尺寸
        width = custom_width if custom_width else config.IMAGE_WIDTH
        height = custom_height if custom_height else config.IMAGE_HEIGHT

        # 確保模型已載入
        initialize_model()

        # 生成前清理 GPU 快取
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
            print("✓ 已清理 GPU 快取")

        # 生成圖像
        print(f"開始生成：{full_prompt}")
        if style_keywords:
            print(f"風格: {style_keywords}")
        # 使用隨機種子,讓每次生成的圖片都不同
        seed = random.randint(0, 2**32 - 1)
        print(f"使用種子: {seed}")

        # 使用配置檔案中的參數生成圖片
        print(f"生成解析度: {width}x{height}")

        # 生成圖片
        # 使用 CUDA generator 以確保與模型在同一設備
        device = "cuda" if torch.cuda.is_available() else "cpu"
        generator = torch.Generator(device=device).manual_seed(seed)

        image = pipe(
            prompt=full_prompt,
            height=height,
            width=width,
            num_inference_steps=config.NUM_INFERENCE_STEPS,
            guidance_scale=config.GUIDANCE_SCALE,
            generator=generator,
        ).images[0]

        # 生成帶有日期時間的檔案名稱
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"generated_{timestamp}.png"
        save_path = os.path.join(output_path, filename)

        # 儲存圖片
        image.save(save_path)
        print(f"圖片已儲存至：{save_path}")

        # 添加到歷史記錄
        add_to_history(prompt, filename)

        # 將圖片轉換為 base64 以便在網頁上顯示
        buffered = BytesIO()
        image.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode()

        return jsonify({
            'success': True,
            'image': f"data:image/png;base64,{img_str}",
            'filename': filename,
            'prompt': prompt,
            'message': f'圖片已成功生成並儲存為 {filename}'
        })

    except Exception as e:
        print(f"錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/images/<filename>')
def get_image(filename):
    """提供圖片下載"""
    return send_from_directory(output_path, filename)

@app.route('/history', methods=['GET'])
def get_history():
    """獲取歷史記錄"""
    try:
        history = load_history()
        return jsonify({
            'success': True,
            'history': history
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/history', methods=['DELETE'])
def clear_history():
    """清除所有歷史記錄"""
    try:
        save_history([])
        return jsonify({
            'success': True,
            'message': '歷史記錄已清除'
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/batch-generate', methods=['POST'])
def batch_generate():
    """批量生成圖片 API"""
    try:
        data = request.get_json()
        prompts = data.get('prompts', [])

        if not prompts or len(prompts) == 0:
            return jsonify({'error': '請輸入至少一個提示詞'}), 400

        # 限制批量數量 (避免 VRAM 問題)
        max_batch = 20
        if len(prompts) > max_batch:
            return jsonify({'error': f'批量生成最多支援 {max_batch} 張圖片'}), 400

        # 確保模型已載入
        initialize_model()

        results = []
        failed_prompts = []

        print(f"\n開始批量生成 {len(prompts)} 張圖片...")

        for idx, prompt in enumerate(prompts, 1):
            prompt = prompt.strip()
            if not prompt:
                continue

            try:
                # 生成前清理 GPU 快取
                if torch.cuda.is_available():
                    torch.cuda.empty_cache()

                print(f"\n[{idx}/{len(prompts)}] 生成：{prompt}")

                # 使用隨機種子
                seed = random.randint(0, 2**32 - 1)
                print(f"使用種子: {seed}")

                # 生成圖片
                device = "cuda" if torch.cuda.is_available() else "cpu"
                generator = torch.Generator(device=device).manual_seed(seed)

                image = pipe(
                    prompt=prompt,
                    height=config.IMAGE_HEIGHT,
                    width=config.IMAGE_WIDTH,
                    num_inference_steps=config.NUM_INFERENCE_STEPS,
                    guidance_scale=config.GUIDANCE_SCALE,
                    generator=generator,
                ).images[0]

                # 生成檔案名稱
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"batch_{timestamp}_{idx:03d}.png"
                save_path = os.path.join(output_path, filename)

                # 儲存圖片
                image.save(save_path)
                print(f"✓ 圖片已儲存: {filename}")

                # 添加到歷史記錄
                add_to_history(prompt, filename)

                # 轉換為 base64
                buffered = BytesIO()
                image.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode()

                results.append({
                    'success': True,
                    'prompt': prompt,
                    'filename': filename,
                    'image': f"data:image/png;base64,{img_str}",
                    'index': idx
                })

            except Exception as e:
                print(f"✗ 生成失敗 [{idx}/{len(prompts)}]: {str(e)}")
                failed_prompts.append({
                    'prompt': prompt,
                    'index': idx,
                    'error': str(e)
                })
                results.append({
                    'success': False,
                    'prompt': prompt,
                    'error': str(e),
                    'index': idx
                })

        print(f"\n批量生成完成! 成功: {len(results) - len(failed_prompts)}/{len(prompts)}")

        return jsonify({
            'success': True,
            'total': len(prompts),
            'succeeded': len(results) - len(failed_prompts),
            'failed': len(failed_prompts),
            'results': results,
            'message': f'批量生成完成，成功 {len(results) - len(failed_prompts)} 張，失敗 {len(failed_prompts)} 張'
        })

    except Exception as e:
        print(f"批量生成錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/batch-download', methods=['POST'])
def batch_download():
    """批量下載圖片為 ZIP"""
    try:
        data = request.get_json()
        filenames = data.get('filenames', [])

        if not filenames:
            return jsonify({'error': '沒有要下載的檔案'}), 400

        # 建立臨時 ZIP 檔案
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        zip_filename = f"batch_images_{timestamp}.zip"
        zip_path = os.path.join(tempfile.gettempdir(), zip_filename)

        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for filename in filenames:
                file_path = os.path.join(output_path, filename)
                if os.path.exists(file_path):
                    zipf.write(file_path, filename)

        # 發送檔案後刪除臨時檔案
        return send_file(
            zip_path,
            mimetype='application/zip',
            as_attachment=True,
            download_name=zip_filename
        )

    except Exception as e:
        print(f"批量下載錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/delete-images', methods=['POST'])
def delete_images():
    """刪除選定的圖片"""
    try:
        data = request.get_json()
        filenames = data.get('filenames', [])

        if not filenames:
            return jsonify({'error': '請選擇要刪除的圖片'}), 400

        deleted_count = 0
        failed_files = []

        # 載入歷史記錄
        history = load_history()

        # 刪除圖片檔案並從歷史記錄中移除
        for filename in filenames:
            file_path = os.path.join(output_path, filename)

            try:
                # 刪除圖片檔案
                if os.path.exists(file_path):
                    os.remove(file_path)
                    deleted_count += 1
                    print(f"✓ 已刪除圖片: {filename}")

                # 從歷史記錄中移除
                history = [item for item in history if item['filename'] != filename]

            except Exception as e:
                print(f"✗ 刪除 {filename} 失敗: {e}")
                failed_files.append(filename)

        # 儲存更新後的歷史記錄
        save_history(history)

        if failed_files:
            return jsonify({
                'success': True,
                'deleted': deleted_count,
                'failed': len(failed_files),
                'failed_files': failed_files,
                'message': f'已刪除 {deleted_count} 張圖片，{len(failed_files)} 張失敗'
            })
        else:
            return jsonify({
                'success': True,
                'deleted': deleted_count,
                'message': f'成功刪除 {deleted_count} 張圖片'
            })

    except Exception as e:
        print(f"刪除圖片錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/templates', methods=['GET'])
def get_templates():
    """獲取風格模板列表"""
    try:
        if os.path.exists(templates_file):
            with open(templates_file, 'r', encoding='utf-8') as f:
                templates = json.load(f)
            return jsonify({
                'success': True,
                'templates': templates
            })
        else:
            return jsonify({
                'success': False,
                'error': '模板檔案不存在'
            }), 404
    except Exception as e:
        print(f"讀取模板錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/size-presets', methods=['GET'])
def get_size_presets():
    """獲取尺寸預設列表"""
    presets = {
        "社群媒體": [
            {"name": "Instagram 正方形", "width": 1080, "height": 1080, "ratio": "1:1"},
            {"name": "Instagram 直式", "width": 1080, "height": 1350, "ratio": "4:5"},
            {"name": "Facebook 封面", "width": 1200, "height": 630, "ratio": "1.91:1"},
            {"name": "Twitter 卡片", "width": 1200, "height": 675, "ratio": "16:9"},
            {"name": "YouTube 縮圖", "width": 1280, "height": 720, "ratio": "16:9"}
        ],
        "列印尺寸": [
            {"name": "A4 直式", "width": 2480, "height": 3508, "ratio": "A4"},
            {"name": "A4 橫式", "width": 3508, "height": 2480, "ratio": "A4"},
            {"name": "A5 直式", "width": 1748, "height": 2480, "ratio": "A5"},
            {"name": "明信片", "width": 1600, "height": 1200, "ratio": "4:3"}
        ],
        "標準尺寸": [
            {"name": "正方形 512", "width": 512, "height": 512, "ratio": "1:1", "vram": "低"},
            {"name": "正方形 768", "width": 768, "height": 768, "ratio": "1:1", "vram": "中"},
            {"name": "正方形 1024", "width": 1024, "height": 1024, "ratio": "1:1", "vram": "高"},
            {"name": "寬屏 16:9", "width": 1024, "height": 576, "ratio": "16:9", "vram": "中"},
            {"name": "直式 9:16", "width": 576, "height": 1024, "ratio": "9:16", "vram": "中"}
        ]
    }

    return jsonify({
        'success': True,
        'presets': presets,
        'current': {
            'width': config.IMAGE_WIDTH,
            'height': config.IMAGE_HEIGHT
        }
    })

@app.route('/tags', methods=['GET'])
def get_all_tags():
    """獲取所有使用過的標籤"""
    try:
        history = load_history()
        all_tags = set()
        tag_counts = {}

        for item in history:
            if 'tags' in item:
                for tag in item['tags']:
                    all_tags.add(tag)
                    tag_counts[tag] = tag_counts.get(tag, 0) + 1

        # 按使用頻率排序
        sorted_tags = sorted(tag_counts.items(), key=lambda x: x[1], reverse=True)

        return jsonify({
            'success': True,
            'tags': [{'name': tag, 'count': count} for tag, count in sorted_tags]
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/history/<item_id>/tags', methods=['POST'])
def update_tags(item_id):
    """更新歷史記錄的標籤"""
    try:
        data = request.get_json()
        tags = data.get('tags', [])

        history = load_history()
        updated = False

        for item in history:
            if item.get('id') == item_id:
                item['tags'] = tags
                updated = True
                break

        if updated:
            save_history(history)
            return jsonify({
                'success': True,
                'message': '標籤已更新'
            })
        else:
            return jsonify({
                'success': False,
                'error': '找不到該記錄'
            }), 404
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/history/filter', methods=['POST'])
def filter_history():
    """根據標籤過濾歷史記錄"""
    try:
        data = request.get_json()
        filter_tags = data.get('tags', [])

        if not filter_tags:
            # 沒有過濾條件，返回全部
            history = load_history()
        else:
            history = load_history()
            # 過濾包含任一標籤的記錄
            filtered = [
                item for item in history
                if 'tags' in item and any(tag in item['tags'] for tag in filter_tags)
            ]
            history = filtered

        return jsonify({
            'success': True,
            'history': history,
            'count': len(history)
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/seed-control', methods=['POST'])
def generate_with_seed():
    """使用指定種子生成圖片（用於重現結果）"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', '')
        seed = data.get('seed')
        style_keywords = data.get('style_keywords', '')
        custom_width = data.get('width')
        custom_height = data.get('height')

        if not prompt:
            return jsonify({'error': '請輸入提示詞'}), 400

        if seed is None:
            seed = random.randint(0, 2**32 - 1)

        # 組合風格關鍵字
        if style_keywords:
            full_prompt = f"{prompt}, {style_keywords}"
        else:
            full_prompt = prompt

        # 確定尺寸
        width = custom_width if custom_width else config.IMAGE_WIDTH
        height = custom_height if custom_height else config.IMAGE_HEIGHT

        # 確保模型已載入
        initialize_model()

        # 生成前清理 GPU 快取
        if torch.cuda.is_available():
            torch.cuda.empty_cache()

        print(f"開始生成（固定種子）：{full_prompt}")
        print(f"種子: {seed}")
        print(f"解析度: {width}x{height}")

        # 生成圖片
        device = "cuda" if torch.cuda.is_available() else "cpu"
        generator = torch.Generator(device=device).manual_seed(seed)

        image = pipe(
            prompt=full_prompt,
            height=height,
            width=width,
            num_inference_steps=config.NUM_INFERENCE_STEPS,
            guidance_scale=config.GUIDANCE_SCALE,
            generator=generator,
        ).images[0]

        # 儲存
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"seed_{seed}_{timestamp}.png"
        save_path = os.path.join(output_path, filename)
        image.save(save_path)

        # 添加到歷史
        add_to_history(prompt, filename)

        # 轉base64
        buffered = BytesIO()
        image.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode()

        return jsonify({
            'success': True,
            'image': f"data:image/png;base64,{img_str}",
            'filename': filename,
            'prompt': prompt,
            'seed': seed,
            'message': f'圖片已生成（種子: {seed}）'
        })
    except Exception as e:
        print(f"錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/add-text-overlay', methods=['POST'])
def add_text_overlay():
    """在圖片上添加文字疊加層"""
    try:
        data = request.get_json()
        filename = data.get('filename', '')
        text = data.get('text', '')
        position = data.get('position', 'bottom')  # top, middle, bottom, custom
        text_color = data.get('text_color', 'white')  # white, black
        bg_overlay = data.get('bg_overlay', True)  # 背景遮罩
        font_size = data.get('font_size', 48)  # 字體大小
        custom_x = data.get('custom_x')  # 自定義 X 座標
        custom_y = data.get('custom_y')  # 自定義 Y 座標

        if not filename:
            return jsonify({'error': '請提供圖片檔名'}), 400

        if not text:
            return jsonify({'error': '請輸入文字內容'}), 400

        # 載入原圖
        image_path = os.path.join(output_path, filename)
        if not os.path.exists(image_path):
            return jsonify({'error': '圖片檔案不存在'}), 404

        image = Image.open(image_path)
        draw = ImageDraw.Draw(image, 'RGBA')

        # 嘗試載入中文字體（Windows 系統）
        try:
            # 常見的 Windows 中文字體路徑
            font_paths = [
                "C:/Windows/Fonts/msyh.ttc",  # 微軟雅黑
                "C:/Windows/Fonts/msjh.ttc",  # 微軟正黑體
                "C:/Windows/Fonts/simsun.ttc",  # 宋體
                "C:/Windows/Fonts/simhei.ttf",  # 黑體
            ]

            font = None
            for font_path in font_paths:
                if os.path.exists(font_path):
                    font = ImageFont.truetype(font_path, font_size)
                    print(f"✓ 載入字體: {font_path}")
                    break

            if font is None:
                # 使用預設字體
                font = ImageFont.load_default()
                print("⚠ 使用預設字體（不支援中文）")
        except Exception as e:
            print(f"字體載入錯誤: {e}")
            font = ImageFont.load_default()

        # 計算文字尺寸
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]

        img_width, img_height = image.size

        # 計算文字位置
        if position == 'custom' and custom_x is not None and custom_y is not None:
            x = custom_x
            y = custom_y
        else:
            # 水平置中
            x = (img_width - text_width) // 2

            # 垂直位置
            if position == 'top':
                y = 50
            elif position == 'middle':
                y = (img_height - text_height) // 2
            else:  # bottom
                y = img_height - text_height - 50

        # 繪製半透明背景遮罩
        if bg_overlay:
            padding = 20
            overlay_color = (0, 0, 0, 180) if text_color == 'white' else (255, 255, 255, 180)
            overlay_bbox = [
                x - padding,
                y - padding,
                x + text_width + padding,
                y + text_height + padding
            ]
            draw.rectangle(overlay_bbox, fill=overlay_color)

        # 繪製文字
        text_rgb = (255, 255, 255) if text_color == 'white' else (0, 0, 0)
        draw.text((x, y), text, font=font, fill=text_rgb)

        # 儲存新圖片
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        new_filename = f"text_overlay_{timestamp}.png"
        new_save_path = os.path.join(output_path, new_filename)
        image.save(new_save_path)
        print(f"✓ 文字疊加圖片已儲存: {new_filename}")

        # 添加到歷史記錄
        add_to_history(f"文字疊加: {text}", new_filename)

        # 轉換為 base64
        buffered = BytesIO()
        image.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode()

        return jsonify({
            'success': True,
            'image': f"data:image/png;base64,{img_str}",
            'filename': new_filename,
            'original_filename': filename,
            'text': text,
            'message': '文字疊加完成'
        })

    except Exception as e:
        print(f"文字疊加錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/export-pdf', methods=['POST'])
def export_pdf():
    """導出多張圖片為 PDF"""
    try:
        data = request.get_json()
        filenames = data.get('filenames', [])
        title = data.get('title', '圖片集')
        include_prompts = data.get('include_prompts', True)
        layout = data.get('layout', 'single')  # single: 一頁一圖, grid: 一頁兩圖

        if not filenames:
            return jsonify({'error': '請選擇至少一張圖片'}), 400

        # 建立臨時 PDF 檔案
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pdf_filename = f"export_{timestamp}.pdf"
        pdf_path = os.path.join(tempfile.gettempdir(), pdf_filename)

        # 建立 PDF
        c = canvas.Canvas(pdf_path, pagesize=A4)
        page_width, page_height = A4

        # 註冊中文字體
        font_registered = False
        try:
            font_paths = [
                "C:/Windows/Fonts/msyh.ttc",  # 微軟雅黑
                "C:/Windows/Fonts/msjh.ttc",  # 微軟正黑體
                "C:/Windows/Fonts/simsun.ttc",  # 宋體
            ]
            for font_path in font_paths:
                if os.path.exists(font_path):
                    pdfmetrics.registerFont(TTFont('ChineseFont', font_path))
                    font_registered = True
                    print(f"✓ PDF 字體已註冊: {font_path}")
                    break
        except Exception as e:
            print(f"⚠ 中文字體註冊失敗: {e}")

        # 載入歷史記錄（用於取得 prompts）
        history = load_history()
        filename_to_prompt = {item['filename']: item['prompt'] for item in history}

        # 繪製封面頁
        c.setFont('ChineseFont' if font_registered else 'Helvetica-Bold', 28)
        c.drawCentredString(page_width / 2, page_height - 2 * inch, title)

        c.setFont('ChineseFont' if font_registered else 'Helvetica', 12)
        c.drawCentredString(page_width / 2, page_height - 2.5 * inch,
                           f"生成日期: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        c.drawCentredString(page_width / 2, page_height - 2.8 * inch,
                           f"共 {len(filenames)} 張圖片")

        # 繪製分隔線
        c.line(100, page_height - 3 * inch, page_width - 100, page_height - 3 * inch)

        c.setFont('ChineseFont' if font_registered else 'Helvetica', 10)
        c.drawCentredString(page_width / 2, 1.5 * inch,
                           "Generated with Z-Image-Turbo Web UI")

        c.showPage()  # 結束封面頁

        # 添加圖片頁
        margin = 0.5 * inch
        usable_width = page_width - 2 * margin
        usable_height = page_height - 2 * margin

        for idx, filename in enumerate(filenames, 1):
            image_path = os.path.join(output_path, filename)

            if not os.path.exists(image_path):
                print(f"⚠ 圖片不存在: {filename}")
                continue

            try:
                # 載入圖片
                img = Image.open(image_path)
                img_width, img_height = img.size

                # 計算縮放比例（保持比例）
                if layout == 'single':
                    # 單圖模式：最大化圖片尺寸
                    max_img_height = usable_height - 1.5 * inch  # 預留標題空間
                    max_img_width = usable_width
                else:
                    # 網格模式：一頁兩圖
                    max_img_height = (usable_height - 2 * inch) / 2
                    max_img_width = usable_width

                scale = min(max_img_width / img_width, max_img_height / img_height)
                scaled_width = img_width * scale
                scaled_height = img_height * scale

                # 繪製標題
                c.setFont('ChineseFont' if font_registered else 'Helvetica-Bold', 14)
                title_text = f"圖片 {idx}/{len(filenames)}"
                c.drawString(margin, page_height - margin - 0.3 * inch, title_text)

                # 繪製提示詞（如果啟用）
                if include_prompts and filename in filename_to_prompt:
                    prompt = filename_to_prompt[filename]
                    c.setFont('ChineseFont' if font_registered else 'Helvetica', 10)

                    # 處理長提示詞（換行）
                    max_width = usable_width
                    words = prompt.split()
                    lines = []
                    current_line = ""

                    for word in words:
                        test_line = current_line + " " + word if current_line else word
                        if c.stringWidth(test_line, 'ChineseFont' if font_registered else 'Helvetica', 10) < max_width:
                            current_line = test_line
                        else:
                            if current_line:
                                lines.append(current_line)
                            current_line = word

                    if current_line:
                        lines.append(current_line)

                    # 限制最多 3 行
                    lines = lines[:3]

                    y_pos = page_height - margin - 0.6 * inch
                    for line in lines:
                        c.drawString(margin, y_pos, line)
                        y_pos -= 0.2 * inch

                # 計算圖片位置（置中）
                img_x = margin + (usable_width - scaled_width) / 2
                img_y = page_height - margin - 1.2 * inch - scaled_height - (0.2 * inch * min(len(lines) if include_prompts and filename in filename_to_prompt else 0, 3))

                # 繪製圖片
                c.drawImage(ImageReader(img), img_x, img_y,
                           width=scaled_width, height=scaled_height,
                           preserveAspectRatio=True)

                # 繪製檔案名稱（底部）
                c.setFont('ChineseFont' if font_registered else 'Helvetica', 8)
                c.drawCentredString(page_width / 2, margin / 2, filename)

                # 繪製頁碼
                c.drawRightString(page_width - margin, margin / 2, f"第 {idx} 頁")

                c.showPage()  # 下一頁

            except Exception as e:
                print(f"✗ 處理圖片 {filename} 時出錯: {e}")
                continue

        # 儲存 PDF
        c.save()
        print(f"✓ PDF 已生成: {pdf_filename}")

        # 發送檔案
        return send_file(
            pdf_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=pdf_filename
        )

    except Exception as e:
        print(f"PDF 導出錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/export-ppt', methods=['POST'])
def export_ppt():
    """導出多張圖片為 PowerPoint"""
    try:
        data = request.get_json()
        filenames = data.get('filenames', [])
        title = data.get('title', '圖片集')
        include_prompts = data.get('include_prompts', True)
        theme = data.get('theme', 'default')  # default, dark, light

        if not filenames:
            return jsonify({'error': '請選擇至少一張圖片'}), 400

        # 建立簡報
        prs = Presentation()
        prs.slide_width = Inches(10)  # 16:9 寬屏
        prs.slide_height = Inches(5.625)

        # 載入歷史記錄
        history = load_history()
        filename_to_prompt = {item['filename']: item['prompt'] for item in history}

        # 添加封面頁
        title_slide_layout = prs.slide_layouts[0]  # 標題投影片
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]

        title_shape.text = title
        subtitle.text = f"生成日期: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n共 {len(filenames)} 張圖片"

        # 設定主題顏色 (使用 RGBColor)
        if theme == 'dark':
            bg_color = RGBColor(30, 30, 30)
            text_color = RGBColor(255, 255, 255)
        elif theme == 'light':
            bg_color = RGBColor(255, 255, 255)
            text_color = RGBColor(0, 0, 0)
        else:
            bg_color = RGBColor(245, 245, 245)
            text_color = RGBColor(50, 50, 50)

        # 添加圖片投影片
        for idx, filename in enumerate(filenames, 1):
            image_path = os.path.join(output_path, filename)

            if not os.path.exists(image_path):
                print(f"⚠ 圖片不存在: {filename}")
                continue

            try:
                # 建立空白投影片
                blank_slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(blank_slide_layout)

                # 設定背景顏色
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

                # 載入圖片以取得尺寸
                img = Image.open(image_path)
                img_width, img_height = img.size

                # 計算圖片在投影片中的尺寸（保持比例）
                max_width = Inches(9)  # 留 0.5 英寸邊距
                max_height = Inches(4.5)

                # 計算縮放比例
                width_scale = max_width / Inches(img_width / 100)
                height_scale = max_height / Inches(img_height / 100)
                scale = min(width_scale, height_scale, 1.0)

                pic_width = Inches(img_width / 100) * scale
                pic_height = Inches(img_height / 100) * scale

                # 圖片置中
                left = (prs.slide_width - pic_width) / 2
                top = Inches(0.5)

                # 添加圖片
                pic = slide.shapes.add_picture(image_path, left, top,
                                              width=pic_width, height=pic_height)

                # 添加標題文字框（如果有提示詞）
                if include_prompts and filename in filename_to_prompt:
                    prompt = filename_to_prompt[filename]

                    # 在底部添加文字框
                    text_box_left = Inches(0.5)
                    text_box_top = top + pic_height + Inches(0.1)
                    text_box_width = Inches(9)
                    text_box_height = Inches(0.8)

                    textbox = slide.shapes.add_textbox(text_box_left, text_box_top,
                                                       text_box_width, text_box_height)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True

                    p = text_frame.paragraphs[0]
                    p.text = prompt[:200]  # 限制長度
                    p.font.size = Pt(12)
                    p.font.color.rgb = text_color
                    p.alignment = PP_ALIGN.CENTER

                # 添加頁碼（右下角）
                page_num_left = Inches(9)
                page_num_top = Inches(5.2)
                page_num_width = Inches(0.8)
                page_num_height = Inches(0.3)

                page_box = slide.shapes.add_textbox(page_num_left, page_num_top,
                                                    page_num_width, page_num_height)
                page_frame = page_box.text_frame
                page_p = page_frame.paragraphs[0]
                page_p.text = f"{idx}/{len(filenames)}"
                page_p.font.size = Pt(10)
                page_p.font.color.rgb = text_color
                page_p.alignment = PP_ALIGN.RIGHT

                print(f"✓ 已添加投影片 {idx}: {filename}")

            except Exception as e:
                print(f"✗ 處理圖片 {filename} 時出錯: {e}")
                continue

        # 儲存 PPT
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_filename = f"export_{timestamp}.pptx"
        ppt_path = os.path.join(tempfile.gettempdir(), ppt_filename)
        prs.save(ppt_path)
        print(f"✓ PPT 已生成: {ppt_filename}")

        # 發送檔案
        return send_file(
            ppt_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=ppt_filename
        )

    except Exception as e:
        print(f"PPT 導出錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # 顯示配置資訊
    config.print_config_info()

    print(f"模型緩存路徑：{cache_path}")
    print(f"生成圖片儲存路徑：{output_path}")

    # 🚀 優化：在伺服器啟動時就預載入模型
    print("\n===================================")
    print("正在預載入模型...")
    print("===================================")
    initialize_model()
    print("✅ 模型已就緒！可以開始生成圖片了\n")

    print("正在啟動 Flask 伺服器...")
    print(f"請在瀏覽器開啟: http://localhost:{config.PORT}")
    print("===================================\n")

    # 關閉 reloader 避免生成過程中重新載入
    app.run(
        debug=config.DEBUG,
        host=config.HOST,
        port=config.PORT,
        use_reloader=config.USE_RELOADER
    )
