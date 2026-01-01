"""
Export Routes - 導出功能相關路由
"""
import os
import tempfile
from datetime import datetime
from flask import Blueprint, request, jsonify, send_file
from PIL import Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import config
from services.history_service import get_history_service


export_bp = Blueprint('export', __name__)


@export_bp.route('/export-pdf', methods=['POST'])
def export_pdf():
    """導出多張圖片為 PDF"""
    try:
        data = request.get_json()
        filenames = data.get('filenames', [])
        title = data.get('title', '圖片集')
        include_prompts = data.get('include_prompts', True)
        layout = data.get('layout', 'single')  # single: 一頁一圖, grid: 一頁兩圖

        if not filenames:
            return jsonify({'error': '請選擇至少一張圖片'}), 400

        # 建立臨時 PDF 檔案
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pdf_filename = f"export_{timestamp}.pdf"
        pdf_path = os.path.join(tempfile.gettempdir(), pdf_filename)

        # 建立 PDF
        c = canvas.Canvas(pdf_path, pagesize=A4)
        page_width, page_height = A4

        # 註冊中文字體
        font_registered = False
        try:
            font_paths = [
                "C:/Windows/Fonts/msyh.ttc",  # 微軟雅黑
                "C:/Windows/Fonts/msjh.ttc",  # 微軟正黑體
                "C:/Windows/Fonts/simsun.ttc",  # 宋體
            ]
            for font_path in font_paths:
                if os.path.exists(font_path):
                    pdfmetrics.registerFont(TTFont('ChineseFont', font_path))
                    font_registered = True
                    print(f"✓ PDF 字體已註冊: {font_path}")
                    break
        except Exception as e:
            print(f"⚠ 中文字體註冊失敗: {e}")

        # 載入歷史記錄（用於取得 prompts）
        history_service = get_history_service()
        history = history_service.load_history()
        filename_to_prompt = {item['filename']: item['prompt'] for item in history}

        # 繪製封面頁
        c.setFont('ChineseFont' if font_registered else 'Helvetica-Bold', 28)
        c.drawCentredString(page_width / 2, page_height - 2 * inch, title)

        c.setFont('ChineseFont' if font_registered else 'Helvetica', 12)
        c.drawCentredString(page_width / 2, page_height - 2.5 * inch,
                           f"生成日期: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        c.drawCentredString(page_width / 2, page_height - 2.8 * inch,
                           f"共 {len(filenames)} 張圖片")

        # 繪製分隔線
        c.line(100, page_height - 3 * inch, page_width - 100, page_height - 3 * inch)

        c.setFont('ChineseFont' if font_registered else 'Helvetica', 10)
        c.drawCentredString(page_width / 2, 1.5 * inch,
                           "Generated with Z-Image-Turbo Web UI")

        c.showPage()  # 結束封面頁

        # 添加圖片頁
        margin = 0.5 * inch
        usable_width = page_width - 2 * margin
        usable_height = page_height - 2 * margin

        for idx, filename in enumerate(filenames, 1):
            image_path = os.path.join(config.OUTPUT_PATH, filename)

            if not os.path.exists(image_path):
                print(f"⚠ 圖片不存在: {filename}")
                continue

            try:
                # 載入圖片
                img = Image.open(image_path)
                img_width, img_height = img.size

                # 計算縮放比例（保持比例）
                if layout == 'single':
                    max_img_height = usable_height - 1.5 * inch
                    max_img_width = usable_width
                else:
                    max_img_height = (usable_height - 2 * inch) / 2
                    max_img_width = usable_width

                scale = min(max_img_width / img_width, max_img_height / img_height)
                scaled_width = img_width * scale
                scaled_height = img_height * scale

                # 繪製標題
                c.setFont('ChineseFont' if font_registered else 'Helvetica-Bold', 14)
                title_text = f"圖片 {idx}/{len(filenames)}"
                c.drawString(margin, page_height - margin - 0.3 * inch, title_text)

                # 繪製提示詞（如果啟用）
                lines = []
                if include_prompts and filename in filename_to_prompt:
                    prompt = filename_to_prompt[filename]
                    c.setFont('ChineseFont' if font_registered else 'Helvetica', 10)

                    max_width = usable_width
                    words = prompt.split()
                    current_line = ""

                    for word in words:
                        test_line = current_line + " " + word if current_line else word
                        if c.stringWidth(test_line, 'ChineseFont' if font_registered else 'Helvetica', 10) < max_width:
                            current_line = test_line
                        else:
                            if current_line:
                                lines.append(current_line)
                            current_line = word

                    if current_line:
                        lines.append(current_line)

                    lines = lines[:3]

                    y_pos = page_height - margin - 0.6 * inch
                    for line in lines:
                        c.drawString(margin, y_pos, line)
                        y_pos -= 0.2 * inch

                # 計算圖片位置（置中）
                img_x = margin + (usable_width - scaled_width) / 2
                img_y = page_height - margin - 1.2 * inch - scaled_height - (0.2 * inch * len(lines))

                # 繪製圖片
                c.drawImage(ImageReader(img), img_x, img_y,
                           width=scaled_width, height=scaled_height,
                           preserveAspectRatio=True)

                # 繪製檔案名稱（底部）
                c.setFont('ChineseFont' if font_registered else 'Helvetica', 8)
                c.drawCentredString(page_width / 2, margin / 2, filename)

                # 繪製頁碼
                c.drawRightString(page_width - margin, margin / 2, f"第 {idx} 頁")

                c.showPage()

            except Exception as e:
                print(f"✗ 處理圖片 {filename} 時出錯: {e}")
                continue

        # 儲存 PDF
        c.save()
        print(f"✓ PDF 已生成: {pdf_filename}")

        return send_file(
            pdf_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name=pdf_filename
        )

    except Exception as e:
        print(f"PDF 導出錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500


@export_bp.route('/export-ppt', methods=['POST'])
def export_ppt():
    """導出多張圖片為 PowerPoint"""
    try:
        data = request.get_json()
        filenames = data.get('filenames', [])
        title = data.get('title', '圖片集')
        include_prompts = data.get('include_prompts', True)
        theme = data.get('theme', 'default')  # default, dark, light

        if not filenames:
            return jsonify({'error': '請選擇至少一張圖片'}), 400

        # 建立簡報
        prs = Presentation()
        prs.slide_width = Inches(10)  # 16:9 寬屏
        prs.slide_height = Inches(5.625)

        # 載入歷史記錄
        history_service = get_history_service()
        history = history_service.load_history()
        filename_to_prompt = {item['filename']: item['prompt'] for item in history}

        # 添加封面頁
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle = slide.placeholders[1]

        title_shape.text = title
        subtitle.text = f"生成日期: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n共 {len(filenames)} 張圖片"

        # 設定主題顏色
        if theme == 'dark':
            bg_color = RGBColor(30, 30, 30)
            text_color = RGBColor(255, 255, 255)
        elif theme == 'light':
            bg_color = RGBColor(255, 255, 255)
            text_color = RGBColor(0, 0, 0)
        else:
            bg_color = RGBColor(245, 245, 245)
            text_color = RGBColor(50, 50, 50)

        # 添加圖片投影片
        for idx, filename in enumerate(filenames, 1):
            image_path = os.path.join(config.OUTPUT_PATH, filename)

            if not os.path.exists(image_path):
                print(f"⚠ 圖片不存在: {filename}")
                continue

            try:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)

                # 設定背景顏色
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = bg_color

                # 載入圖片以取得尺寸
                img = Image.open(image_path)
                img_width, img_height = img.size

                # 計算圖片在投影片中的尺寸
                max_width = Inches(9)
                max_height = Inches(4.5)

                width_scale = max_width / Inches(img_width / 100)
                height_scale = max_height / Inches(img_height / 100)
                scale = min(width_scale, height_scale, 1.0)

                pic_width = Inches(img_width / 100) * scale
                pic_height = Inches(img_height / 100) * scale

                left = (prs.slide_width - pic_width) / 2
                top = Inches(0.5)

                slide.shapes.add_picture(image_path, left, top,
                                        width=pic_width, height=pic_height)

                # 添加標題文字框
                if include_prompts and filename in filename_to_prompt:
                    prompt = filename_to_prompt[filename]

                    text_box_left = Inches(0.5)
                    text_box_top = top + pic_height + Inches(0.1)
                    text_box_width = Inches(9)
                    text_box_height = Inches(0.8)

                    textbox = slide.shapes.add_textbox(text_box_left, text_box_top,
                                                       text_box_width, text_box_height)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True

                    p = text_frame.paragraphs[0]
                    p.text = prompt[:200]
                    p.font.size = Pt(12)
                    p.font.color.rgb = text_color
                    p.alignment = PP_ALIGN.CENTER

                # 添加頁碼
                page_num_left = Inches(9)
                page_num_top = Inches(5.2)
                page_num_width = Inches(0.8)
                page_num_height = Inches(0.3)

                page_box = slide.shapes.add_textbox(page_num_left, page_num_top,
                                                    page_num_width, page_num_height)
                page_frame = page_box.text_frame
                page_p = page_frame.paragraphs[0]
                page_p.text = f"{idx}/{len(filenames)}"
                page_p.font.size = Pt(10)
                page_p.font.color.rgb = text_color
                page_p.alignment = PP_ALIGN.RIGHT

                print(f"✓ 已添加投影片 {idx}: {filename}")

            except Exception as e:
                print(f"✗ 處理圖片 {filename} 時出錯: {e}")
                continue

        # 儲存 PPT
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_filename = f"export_{timestamp}.pptx"
        ppt_path = os.path.join(tempfile.gettempdir(), ppt_filename)
        prs.save(ppt_path)
        print(f"✓ PPT 已生成: {ppt_filename}")

        return send_file(
            ppt_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=ppt_filename
        )

    except Exception as e:
        print(f"PPT 導出錯誤：{str(e)}")
        return jsonify({'error': str(e)}), 500
